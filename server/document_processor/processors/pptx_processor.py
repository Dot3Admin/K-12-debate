import os
from typing import Dict, List
import logging

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.error("python-pptx not available")

from . import BaseProcessor


class PPTXProcessor(BaseProcessor):
    def process(self, file_path: str) -> Dict:
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed")
        
        logger.info(f"Processing PPTX: {file_path}")
        
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'formulas': [],
            'metadata': {}
        }
        
        try:
            prs = Presentation(file_path)
            
            result['text'] = self._extract_text(prs)
            result['tables'] = self._extract_tables(prs)
            result['metadata'] = self._extract_metadata(prs, file_path)
            
        except Exception as e:
            logger.error(f"PPTX processing error: {e}")
            raise
        
        logger.info(f"PPTX processing complete: {len(result['text'])} chars, {len(result['tables'])} tables")
        
        return result
    
    def _extract_text(self, prs) -> str:
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"# Slide {slide_num}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    
                    if text:
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    level = paragraph.level if hasattr(paragraph, 'level') else 0
                                    indent = "  " * level
                                    text_parts.append(f"{indent}- {para_text}")
                        else:
                            text_parts.append(f"- {text}")
            
            text_parts.append("")
        
        return '\n'.join(text_parts)
    
    def _extract_tables(self, prs) -> List[Dict]:
        tables = []
        table_index = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    
                    table_dict = {
                        'id': f'table_{table_index:03d}',
                        'slide': slide_num,
                        'data': [],
                        'shape': (len(table.rows), len(table.columns))
                    }
                    
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        
                        for cell in row.cells:
                            cell_text = ''.join([p.text for p in cell.text_frame.paragraphs])
                            row_data.append(cell_text)
                        
                        table_dict['data'].append(row_data)
                    
                    tables.append(table_dict)
                    table_index += 1
        
        logger.info(f"Extracted {len(tables)} tables")
        return tables
    
    def _extract_metadata(self, prs, file_path: str) -> Dict:
        metadata = {
            'file_size': os.path.getsize(file_path),
            'file_name': os.path.basename(file_path),
            'total_slides': len(prs.slides),
        }
        
        if hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height'):
            metadata['slide_width'] = prs.slide_width
            metadata['slide_height'] = prs.slide_height
        
        try:
            core_props = prs.core_properties
            
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.subject:
                metadata['subject'] = core_props.subject
            if core_props.created:
                metadata['created'] = str(core_props.created)
            if core_props.modified:
                metadata['modified'] = str(core_props.modified)
                
        except Exception as e:
            logger.warning(f"Properties extraction failed: {e}")
        
        return metadata
