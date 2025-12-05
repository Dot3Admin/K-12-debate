#!/usr/bin/env python3
"""
PPT/PPTX에서 embedded 이미지를 직접 추출하는 스크립트
- Nested shapes (그룹 내 이미지) 지원
- Chart, SmartArt 내 이미지 지원
- 작은 이미지도 포함 (50x50 픽셀 이상)
"""
import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_image_from_shape(shape, slide_num, shape_idx, output_dir, image_counter, slide):
    """
    단일 shape에서 이미지를 추출하는 헬퍼 함수
    
    Returns:
        dict or None: 이미지 메타데이터 또는 None
    """
    try:
        if not hasattr(shape, 'image'):
            return None
            
        # 이미지 데이터 추출
        image = shape.image
        image_bytes = image.blob
        
        # PIL로 이미지 로드
        pil_image = Image.open(io.BytesIO(image_bytes))
        
        # 너무 작은 이미지는 무시 (아이콘, 로고 등) - 기준을 50x50으로 완화
        if pil_image.width < 50 or pil_image.height < 50:
            return None
        
        # 확장자 결정
        image_ext = image.ext
        if not image_ext:
            image_ext = "png"
        
        # 캡션 찾기 (슬라이드에서 텍스트 박스 검색)
        caption = ""
        image_left = shape.left
        image_top = shape.top
        image_width = shape.width
        image_height = shape.height
        
        # 이미지 주변의 텍스트 박스 찾기
        for text_shape in slide.shapes:
            if hasattr(text_shape, 'has_text_frame') and text_shape.has_text_frame:
                text_left = text_shape.left
                text_top = text_shape.top
                
                # 이미지 바로 아래나 위에 있는 텍스트
                is_below = (text_top > image_top + image_height and 
                           text_top < image_top + image_height + 500000)  # EMU units
                is_above = (text_top + text_shape.height < image_top and
                           text_top + text_shape.height > image_top - 500000)
                
                if is_below or is_above:
                    text_content = text_shape.text.strip()
                    if text_content and len(text_content) < 200:  # 캡션은 보통 짧음
                        caption += " " + text_content
        
        caption = caption.strip()
        
        # 이미지 저장
        image_filename = f"img_{image_counter:03d}_slide{slide_num+1}_{shape_idx}.{image_ext}"
        image_path = Path(output_dir) / image_filename
        
        with open(image_path, "wb") as img_file:
            img_file.write(image_bytes)
        
        return {
            "slide": slide_num + 1,
            "image_index": shape_idx,
            "path": str(image_path),
            "caption": caption,
            "width": pil_image.width,
            "height": pil_image.height
        }
        
    except Exception as e:
        print(f"Error extracting image: {e}", file=sys.stderr)
        return None

def process_shapes_recursive(shapes, slide_num, output_dir, image_counter, extracted_images, slide):
    """
    재귀적으로 shapes를 탐색하여 이미지 추출 (nested groups 지원)
    """
    for shape_idx, shape in enumerate(shapes):
        try:
            # 디버그: shape type 출력
            shape_type_name = "UNKNOWN"
            try:
                shape_type_name = str(shape.shape_type)
            except:
                pass
            
            print(f"[DEBUG] Slide {slide_num+1}, Shape {shape_idx}: type={shape_type_name}", file=sys.stderr)
            
            # PICTURE type 이미지
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"[DEBUG] Found PICTURE on slide {slide_num+1}", file=sys.stderr)
                result = extract_image_from_shape(shape, slide_num, shape_idx, output_dir, image_counter, slide)
                if result:
                    extracted_images.append(result)
                    image_counter += 1
            
            # GROUP type - 재귀적으로 탐색
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print(f"[DEBUG] Found GROUP on slide {slide_num+1}, recursing...", file=sys.stderr)
                if hasattr(shape, 'shapes'):
                    image_counter = process_shapes_recursive(
                        shape.shapes, slide_num, output_dir, image_counter, extracted_images, slide
                    )
            
            # PLACEHOLDER - 플레이스홀더 내부 확인
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                print(f"[DEBUG] Found PLACEHOLDER on slide {slide_num+1}", file=sys.stderr)
                # 플레이스홀더 안에도 이미지가 있을 수 있음
                if hasattr(shape, 'image'):
                    result = extract_image_from_shape(shape, slide_num, shape_idx, output_dir, image_counter, slide)
                    if result:
                        extracted_images.append(result)
                        image_counter += 1
            
            # 모든 shape에 대해 이미지 속성 확인 (타입에 상관없이)
            if hasattr(shape, 'image') and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                print(f"[DEBUG] Shape type {shape_type_name} has image attribute!", file=sys.stderr)
                result = extract_image_from_shape(shape, slide_num, shape_idx, output_dir, image_counter, slide)
                if result:
                    extracted_images.append(result)
                    image_counter += 1
            
        except Exception as e:
            print(f"[DEBUG] Error processing shape {shape_idx} on slide {slide_num}: {e}", file=sys.stderr)
            continue
    
    return image_counter

def extract_images_from_pptx(pptx_path, output_dir):
    """
    PPTX에서 모든 이미지를 추출하고 메타데이터 반환 (재귀적 탐색 포함)
    
    Args:
        pptx_path: PPTX 파일 경로
        output_dir: 이미지 저장 디렉토리
    
    Returns:
        list: [{slide, image_index, path, caption, width, height}, ...]
    """
    try:
        prs = Presentation(pptx_path)
        extracted_images = []
        image_counter = 0
        
        for slide_num, slide in enumerate(prs.slides):
            # 재귀적으로 모든 shapes 처리
            image_counter = process_shapes_recursive(
                slide.shapes, slide_num, output_dir, image_counter, extracted_images, slide
            )
        
        return extracted_images
        
    except Exception as e:
        print(f"Error processing PPTX: {e}", file=sys.stderr)
        return []


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(json.dumps({"error": "Usage: extract_images_from_pptx.py <pptx_path> <output_dir>"}))
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]
    
    # 출력 디렉토리 생성
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # 이미지 추출
    images = extract_images_from_pptx(pptx_path, output_dir)
    
    # 결과 출력
    print(json.dumps({
        "success": True,
        "images": images,
        "total": len(images)
    }))
